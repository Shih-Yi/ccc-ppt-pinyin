"""Add Hanyu Pinyin below Chinese lyric lines in a PPTX, each syllable
centered under its character (ruby-style).

Any pre-existing pinyin-only lines (e.g. pinyin sitting ABOVE the lyric
line in the source deck) are removed first, so pinyin only ever appears
below the lyrics — and re-running the tool never duplicates lines.

Alignment method: the pinyin line is left-aligned and padded with spaces.
Positions are computed from font metrics (Arial == Liberation Sans metrics),
so the result renders identically in PowerPoint, LibreOffice and Keynote —
no reliance on custom tab stops, which some renderers ignore.

Pure function: add_pinyin(src) -> BytesIO. Input is not mutated.
"""
import copy
import io
import os
import re
from functools import lru_cache

from fontTools.ttLib import TTFont
from pptx import Presentation
from pptx.util import Emu
from pypinyin import pinyin, Style, load_single_dict, load_phrases_dict

# 教會歌詞常見讀音修正(可自行擴充)
CHAR_OVERRIDES = {
    "祢": "nǐ",   # 對神的第二人稱,字典音為 mí(姓氏)
    "袮": "nǐ",
}
PHRASE_OVERRIDES = {
    "尊主為大": ["zūn", "zhǔ", "wéi", "dà"],
    "尊主为大": ["zūn", "zhǔ", "wéi", "dà"],
}
load_single_dict({ord(c): p for c, p in CHAR_OVERRIDES.items()})
load_phrases_dict({k: [[s] for s in v] for k, v in PHRASE_OVERRIDES.items()})

CJK_RE = re.compile(r"[\u4e00-\u9fff\u3400-\u4dbf]")
# \u8072\u8abf\u7b26\u865f:\u5224\u65b7\u4e00\u884c\u662f\u5426\u70ba\u65e2\u6709\u62fc\u97f3\u884c\u7684\u5f37\u8a0a\u865f
PINYIN_TONE_RE = re.compile(r"[\u0101\u00e1\u01ce\u00e0\u0113\u00e9\u011b\u00e8\u012b\u00ed\u01d0\u00ec\u014d\u00f3\u01d2\u00f2\u016b\u00fa\u01d4\u00f9\u01d6\u01d8\u01da\u01dc\u0144\u0148\u01f9\u1e3f]")
# \u62fc\u97f3\u97f3\u7bc0\u5141\u8a31\u7684\u5b57\u5143(\u5b57\u6bcd\u3001\u8072\u8abf\u6bcd\u97f3\u3001\u00fc\u3001\u9694\u97f3\u865f\u3001\u9023\u5b57\u865f)
PINYIN_TOKEN_RE = re.compile(
    r"^[A-Za-z\u00fc\u00dc\u0101\u00e1\u01ce\u00e0\u0113\u00e9\u011b\u00e8\u012b\u00ed\u01d0\u00ec\u014d\u00f3\u01d2\u00f2\u016b\u00fa\u01d4\u00f9\u01d6\u01d8\u01da\u01dc\u0144\u0148\u01f9\u1e3f'\u2019\-\u00b7]+[,,.\u3002!!??::;;]?$"
)
# bump when output logic changes, so cached results upstream are invalidated
PINYIN_VERSION = 2

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
EMU_PER_PT = 12700
_TTF_CANDIDATES = (
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",  # Streamlit Cloud
    "/Library/Fonts/Arial.ttf",                                          # macOS
    "/System/Library/Fonts/Supplemental/Arial.ttf",                      # macOS
)
SANS_TTF = next((p for p in _TTF_CANDIDATES if os.path.exists(p)), _TTF_CANDIDATES[0])


@lru_cache(maxsize=1)
def _metrics():
    f = TTFont(SANS_TTF)
    upm = f["head"].unitsPerEm
    cmap = f.getBestCmap()
    hmtx = f["hmtx"]
    return upm, cmap, hmtx


def text_width_em(s: str) -> float:
    """Width of `s` in em units, using Arial/Liberation Sans metrics."""
    upm, cmap, hmtx = _metrics()
    w = 0
    for ch in s:
        glyph = cmap.get(ord(ch))
        if glyph is None:  # strip combining marks, retry base char
            import unicodedata
            base = unicodedata.normalize("NFD", ch)[0]
            glyph = cmap.get(ord(base), cmap.get(ord(" ")))
        w += hmtx[glyph][0]
    return w / upm


SPACE_EM = None  # filled lazily


def para_text(p):
    return "".join(r.text for r in p.runs)


def is_pinyin_line(text: str) -> bool:
    """True if `text` looks like a hanyu-pinyin-only line (an existing
    pinyin annotation), e.g. "zūn zhǔ wéi dà". Requires at least one tone
    mark so plain English lyric lines are never touched."""
    t = text.strip()
    if not t or CJK_RE.search(t):
        return False
    if not PINYIN_TONE_RE.search(t):
        return False
    return all(PINYIN_TOKEN_RE.match(tok) for tok in t.split())


def strip_pinyin_paragraphs(tf):
    """Remove every existing pinyin-only paragraph from a text frame
    (wherever it sits — above or below the lyric line), so pinyin can be
    re-inserted cleanly below each lyric. Keeps at least one paragraph so
    the txBody stays schema-valid."""
    paras = list(tf.paragraphs)
    removable = [p for p in paras if is_pinyin_line(para_text(p))]
    if len(removable) == len(paras) and removable:
        # a pinyin-only text frame: keep one (blanked) paragraph so the
        # txBody stays schema-valid
        keep = removable.pop()
        for r in list(keep._p.findall(_qn("r"))):
            keep._p.remove(r)
    for p in removable:
        p._p.getparent().remove(p._p)


def para_font_pt(p):
    for r in p.runs:
        if r.font.size is not None:
            return r.font.size.pt
    return None


def char_cells(text, em_pt):
    """Per-character layout cells: (width_pt, syllable_or_None)."""
    # errors=list -> every non-han char becomes its own item, so the
    # result maps 1:1 onto the characters of `text`
    syls = pinyin(text, style=Style.TONE, errors=lambda chunk: list(chunk))
    cells = []
    for ch, syl in zip(text, syls):
        if CJK_RE.match(ch):
            cells.append((em_pt, syl[0]))
        elif ch in ("\u3000", "\t"):
            cells.append((em_pt, None))
        else:
            cells.append((em_pt / 2, None))  # ASCII space ≈ 0.5em in CJK fonts
    return cells


def _solve_lefts(ideal_lefts, widths, gap):
    """Minimum-deviation left edges subject to no-overlap (L[i+1] >= L[i]+w[i]+gap).
    Transform to a monotone problem and solve with pool-adjacent-violators, so
    colliding syllables share the shift symmetrically instead of drifting right."""
    offsets = [0.0]
    for w in widths[:-1]:
        offsets.append(offsets[-1] + w + gap)
    adj = [l - o for l, o in zip(ideal_lefts, offsets)]  # want non-decreasing
    # PAV: pool blocks whose means violate monotonicity
    blocks = []  # (mean, count)
    for a in adj:
        blocks.append((a, 1))
        while len(blocks) > 1 and blocks[-2][0] > blocks[-1][0]:
            m2, c2 = blocks.pop()
            m1, c1 = blocks.pop()
            blocks.append(((m1 * c1 + m2 * c2) / (c1 + c2), c1 + c2))
    fitted = []
    for m, c in blocks:
        fitted.extend([m] * c)
    return [f + o for f, o in zip(fitted, offsets)]


def build_padded_pinyin(cells, area_w_pt, py_pt):
    """Left-aligned string of spaces + syllables, each syllable centered
    under its character cell where physically possible; colliding wide
    syllables are spread symmetrically. Cumulative targets -> no drift."""
    space_w = text_width_em(" ") * py_pt
    line_w = sum(w for w, _ in cells)
    left = max(0.0, (area_w_pt - line_w) / 2)

    ideal_lefts, widths, syls = [], [], []
    x = left
    for w, syl in cells:
        if syl is not None:
            syl_w = text_width_em(syl) * py_pt
            ideal_lefts.append(x + w / 2 - syl_w / 2)
            widths.append(syl_w)
            syls.append(syl)
        x += w
    lefts = _solve_lefts(ideal_lefts, widths, space_w)

    out = []
    cur = 0.0
    for target, syl_w, syl in zip(lefts, widths, syls):
        n = round((target - cur) / space_w)
        if out and n < 1:
            n = 1  # never let two syllables touch
        n = max(n, 0)
        out.append(" " * n + syl)
        cur += n * space_w + syl_w
    return "".join(out)


def _qn(tag):
    return f"{{{A_NS}}}{tag}"


def add_pinyin(src, min_pt: float = 40, pinyin_pt: float = 20, latin_font: str = "Arial"):
    """Insert a syllable-aligned pinyin line below every Chinese paragraph
    whose font size >= min_pt. pinyin_pt is the absolute pinyin font size in
    points. Returns io.BytesIO of the new .pptx."""
    prs = Presentation(src)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            strip_pinyin_paragraphs(tf)  # 先拿掉既有拼音行(通常在歌詞上面)
            l_ins = tf.margin_left if tf.margin_left is not None else Emu(91440)
            r_ins = tf.margin_right if tf.margin_right is not None else Emu(91440)
            area_w_pt = (int(shape.width) - int(l_ins) - int(r_ins)) / EMU_PER_PT

            for p in list(tf.paragraphs):
                text = para_text(p)
                if not CJK_RE.search(text):
                    continue
                size = para_font_pt(p)
                if size is None or size < min_pt:
                    continue

                py_pt = pinyin_pt
                cells = char_cells(text, size)
                py_line = build_padded_pinyin(cells, area_w_pt, py_pt)

                src_p = p._p
                new_p = copy.deepcopy(src_p)

                pPr = new_p.find(_qn("pPr"))
                if pPr is None:
                    pPr = new_p.makeelement(_qn("pPr"), {})
                    new_p.insert(0, pPr)
                pPr.set("algn", "l")
                pPr.set("marL", "0")
                pPr.set("indent", "0")

                runs = new_p.findall(_qn("r"))
                for extra in runs[1:]:
                    new_p.remove(extra)
                run = runs[0]
                t = run.find(_qn("t"))
                t.text = py_line
                t.set("{http://www.w3.org/XML/1998/namespace}space", "preserve")

                rPr = run.find(_qn("rPr"))
                if rPr is None:
                    rPr = run.makeelement(_qn("rPr"), {})
                    run.insert(0, rPr)
                rPr.set("sz", str(int(py_pt * 100)))
                rPr.set("i", "1")  # 拼音用斜體
                for tag in ("latin", "ea", "cs"):
                    el = rPr.find(_qn(tag))
                    if el is None:
                        el = rPr.makeelement(_qn(tag), {})
                        rPr.append(el)
                    el.set("typeface", latin_font)

                src_p.addnext(new_p)  # 拼音在中文行的正下方

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    return out


if __name__ == "__main__":
    import sys
    src, dst = sys.argv[1], sys.argv[2]
    data = add_pinyin(src)
    with open(dst, "wb") as f:
        f.write(data.read())
    print("written:", dst)
