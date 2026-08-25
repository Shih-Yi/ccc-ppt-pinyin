"""Tests for pinyin_pptx: existing-pinyin removal + insertion below lyrics."""
import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from pinyin_pptx import CJK_RE, DEFAULT_GAP_PCT, add_pinyin, is_pinyin_line


def make_pptx(lines, size_pt=40):
    """One slide, one textbox; each item in `lines` becomes a paragraph."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qn(tag):
    return f"{{{A_NS}}}{tag}"


def _qn_p(tag):
    return f"{{{P_NS}}}{tag}"


def make_pptx_size_in_liststyle(lines, size_pt=48):
    """Like make_pptx, but the runs carry no `sz` at all — the size lives in
    the shape's lstStyle, the way Google Slides exports write it."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.add_run().text = line

    tx_body = tf._txBody
    lst = tx_body.makeelement(_qn("lstStyle"), {})
    lvl1 = lst.makeelement(_qn("lvl1pPr"), {})
    def_rpr = lst.makeelement(_qn("defRPr"), {"sz": str(int(size_pt * 100))})
    lvl1.append(def_rpr)
    lst.append(lvl1)
    tx_body.insert(1, lst)  # right after bodyPr

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def make_pptx_grouped(lines, size_pt=48, scale=1.0):
    """The lyric text box sits inside a group shape. `scale` sets the group's
    ext/chExt ratio, i.e. how child coordinates map onto slide units."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
    tf = box.text_frame
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size_pt)

    xfrm = group._element.find(f"{_qn_p('grpSpPr')}/{_qn('xfrm')}")
    ch_ext = xfrm.find(_qn("chExt"))
    ext = xfrm.find(_qn("ext"))
    ext.set("cx", str(int(int(ch_ext.get("cx")) * scale)))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def all_para_texts(buf):
    """Every paragraph in reading order, descending into groups — read back
    independently of the walker under test."""
    prs = Presentation(buf)
    texts = []

    def visit(container):
        for shape in container.shapes:
            if shape._element.tag == _qn_p("grpSp"):
                visit(shape)
            elif shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    texts.append("".join(r.text for r in p.runs))

    for slide in prs.slides:
        visit(slide)
    return texts


class TestIsPinyinLine:
    def test_detects_tone_marked_pinyin(self):
        assert is_pinyin_line("zūn zhǔ wéi dà")
        assert is_pinyin_line("  yē sū ài nǐ  ")

    def test_rejects_chinese(self):
        assert not is_pinyin_line("尊主為大")
        assert not is_pinyin_line("zūn 主")

    def test_rejects_plain_english(self):
        assert not is_pinyin_line("Jesus loves me")
        assert not is_pinyin_line("Amazing grace")

    def test_rejects_empty(self):
        assert not is_pinyin_line("")
        assert not is_pinyin_line("   ")


class TestAddPinyin:
    def test_pinyin_added_below_lyric(self):
        out = add_pinyin(make_pptx(["尊主為大"]))
        texts = all_para_texts(out)
        assert texts[0] == "尊主為大"
        assert "zūn" in texts[1] and "dà" in texts[1]

    def test_existing_pinyin_above_is_removed(self):
        out = add_pinyin(make_pptx(["zūn zhǔ wéi dà", "尊主為大"]))
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts[0] == "尊主為大"          # 拼音不再出現在歌詞上面
        assert "zūn" in texts[1]               # 新拼音在歌詞下面
        assert len(texts) == 2

    def test_idempotent_rerun(self):
        once = add_pinyin(make_pptx(["尊主為大"]))
        twice = add_pinyin(once)
        texts = [t for t in all_para_texts(twice) if t.strip()]
        assert len(texts) == 2

    def test_english_lines_untouched(self):
        out = add_pinyin(make_pptx(["Jesus loves me", "尊主為大"]))
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts[0] == "Jesus loves me"
        assert texts[1] == "尊主為大"

    def test_small_font_skipped(self):
        out = add_pinyin(make_pptx(["尊主為大"], size_pt=20), min_pt=40)
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts == ["尊主為大"]

    def test_pinyin_only_textbox_is_blanked(self):
        out = add_pinyin(make_pptx(["zūn zhǔ wéi dà"]))
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts == []


class TestInheritedFontSize:
    """Runs without an explicit `sz` must not be skipped — the size is
    inherited from the shape's lstStyle (Google Slides export shape)."""

    def test_size_from_liststyle_is_honoured(self):
        out = add_pinyin(make_pptx_size_in_liststyle(["尊主為大"], size_pt=48))
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts[0] == "尊主為大"
        assert "zūn" in texts[1] and "dà" in texts[1]

    def test_inherited_size_below_threshold_still_skipped(self):
        out = add_pinyin(make_pptx_size_in_liststyle(["尊主為大"], size_pt=28),
                         min_pt=40)
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts == ["尊主為大"]

    def test_explicit_run_size_wins_over_liststyle(self):
        """A 28pt run inside a 48pt shape stays untouched at min_pt=40."""
        buf = make_pptx_size_in_liststyle(["尊主為大"], size_pt=48)
        prs = Presentation(buf)
        run = prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(28)
        patched = io.BytesIO()
        prs.save(patched)
        patched.seek(0)

        texts = [t for t in all_para_texts(add_pinyin(patched, min_pt=40))
                 if t.strip()]
        assert texts == ["尊主為大"]


class TestTrailingWhitespace:
    """A trailing space must not shift the pinyin row: renderers drop it when
    centering, so counting it would offset every syllable to the left."""

    def test_trailing_space_does_not_shift_pinyin(self):
        clean = all_para_texts(add_pinyin(make_pptx(["尊主為大"])))[1]
        spaced = all_para_texts(add_pinyin(make_pptx(["尊主為大 "])))[1]
        assert spaced == clean

    def test_trailing_ideographic_space_does_not_shift_pinyin(self):
        clean = all_para_texts(add_pinyin(make_pptx(["尊主為大"])))[1]
        spaced = all_para_texts(add_pinyin(make_pptx(["尊主為大　"])))[1]
        assert spaced == clean

    def test_leading_space_still_counts(self):
        """Leading whitespace is drawn, so it must still move the pinyin."""
        clean = all_para_texts(add_pinyin(make_pptx(["尊主為大"])))[1]
        led = all_para_texts(add_pinyin(make_pptx([" 尊主為大"])))[1]
        assert len(led) - len(led.lstrip()) > len(clean) - len(clean.lstrip())


class TestReadingOverrides:
    """Readings the dictionary gets wrong for sung worship lyrics."""

    @pytest.mark.parametrize("line, expected", [
        ("祢是我的神", "nǐ shì wǒ di shén"),
        ("神的兒子", "shén di ér zi"),
        ("尊主為大", "zūn zhǔ wéi dà"),
    ])
    def test_overridden_readings(self, line, expected):
        py = all_para_texts(add_pinyin(make_pptx([line])))[1]
        assert py.split() == expected.split()


def _paragraph_elements(buf, keep):
    """The <a:p> element of every paragraph whose text satisfies `keep`."""
    prs = Presentation(buf)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(r.text for r in para.runs)
                if keep(text):
                    out.append(para._p)
    return out


def pinyin_paragraph_elements(buf):
    """The <a:p> element of every pinyin line in the deck."""
    return _paragraph_elements(buf, is_pinyin_line)


def lyric_paragraph_elements(buf):
    """The <a:p> element of every Chinese lyric line in the deck."""
    return _paragraph_elements(buf, lambda t: bool(CJK_RE.search(t)))


def _line_height(p):
    """The paragraph's explicit line height in per-mille, or None."""
    pPr = p.find(_qn("pPr"))
    lnSpc = None if pPr is None else pPr.find(_qn("lnSpc"))
    return None if lnSpc is None else lnSpc.find(_qn("spcPct")).get("val")


class TestLyricToPinyinGap:
    """A paragraph's line height sets the distance down to the line that
    FOLLOWS it. So the lyric-to-pinyin gap lives on the lyric paragraph; put
    it on the pinyin paragraph instead and it squeezes whatever comes after
    the pinyin — the next lyric line, or the English line — which must not
    move."""

    def test_the_gap_is_written_on_the_lyric_paragraph(self):
        p = lyric_paragraph_elements(
            add_pinyin(make_pptx(["尊主為大"]), gap_pct=45))[0]
        assert _line_height(p) == "45000"

    def test_the_pinyin_paragraph_never_gets_its_own_line_height(self):
        """Its line height is the gap BELOW the pinyin, which stays as the
        deck had it."""
        p = pinyin_paragraph_elements(
            add_pinyin(make_pptx(["尊主為大"]), gap_pct=45))[0]
        assert _line_height(p) is None

    def test_the_default_gap_is_written_on_the_lyric_paragraph(self):
        out = add_pinyin(make_pptx(["尊主為大"]))
        assert _line_height(lyric_paragraph_elements(out)[0]) == \
            str(DEFAULT_GAP_PCT * 1000)
        assert _line_height(pinyin_paragraph_elements(out)[0]) is None

    def test_zero_leaves_the_decks_own_line_spacing_alone(self):
        """The escape hatch for a deck whose spacing is already right."""
        out = add_pinyin(make_pptx(["尊主為大"]), gap_pct=0)
        assert _line_height(lyric_paragraph_elements(out)[0]) is None
        assert _line_height(pinyin_paragraph_elements(out)[0]) is None

    def test_no_space_is_added_before_the_pinyin_line(self):
        """Inherited space-before would sit inside the lyric-to-pinyin gap."""
        p = pinyin_paragraph_elements(add_pinyin(make_pptx(["尊主為大"])))[0]
        spc_bef = p.find(_qn("pPr")).find(_qn("spcBef"))
        assert spc_bef.find(_qn("spcPts")).get("val") == "0"

    def test_spacing_elements_lead_pPr(self):
        """DrawingML requires lnSpc, then spcBef, before every other child —
        PowerPoint rejects the file otherwise."""
        out = add_pinyin(make_pptx(["尊主為大"]), gap_pct=45)
        lyric = [c.tag for c in lyric_paragraph_elements(out)[0].find(_qn("pPr"))]
        assert lyric[0] == _qn("lnSpc")
        py = [c.tag for c in pinyin_paragraph_elements(out)[0].find(_qn("pPr"))]
        assert py[0] == _qn("spcBef")


class TestGroupedShapes:
    """Text boxes nested inside a group must be processed too — a group has
    no text frame of its own, so a top-level-only loop skips everything."""

    def test_lyrics_inside_a_group_get_pinyin(self):
        out = add_pinyin(make_pptx_grouped(["尊主為大"]))
        texts = [t for t in all_para_texts(out) if t.strip()]
        assert texts[0] == "尊主為大"
        assert "zūn" in texts[1] and "dà" in texts[1]

    def test_group_child_coordinates_are_scaled(self):
        """A group scaled 2x makes its child twice as wide on the slide, so
        the pinyin gets centered with more leading padding."""
        plain = add_pinyin(make_pptx_grouped(["尊主為大"], scale=1.0))
        wide = add_pinyin(make_pptx_grouped(["尊主為大"], scale=2.0))
        lead = lambda buf: len(all_para_texts(buf)[1]) - len(
            all_para_texts(buf)[1].lstrip())
        assert lead(wide) > lead(plain)

    def test_nested_groups_are_reached(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        inner = slide.shapes.add_group_shape().shapes.add_group_shape()
        box = inner.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        run = box.text_frame.paragraphs[0].add_run()
        run.text = "尊主為大"
        run.font.size = Pt(48)
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        texts = [t for t in all_para_texts(add_pinyin(buf)) if t.strip()]
        assert texts[0] == "尊主為大"
        assert "zūn" in texts[1]
